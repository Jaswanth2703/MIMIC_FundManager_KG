"""
Thesis PPT Generator — 45-60 Minute Defense Presentation
=========================================================
Uses python-pptx to create a professional M.Tech thesis defense PPT.
All citations in APA format. No fake references.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
DIAG_DIR = os.path.join(SCRIPT_DIR, 'diagrams')
OUT_PPT = os.path.join(SCRIPT_DIR, 'thesis_defense_presentation.pptx')

# Colors
DARK_BLUE = RGBColor(0x0D, 0x47, 0xA1)
MED_BLUE = RGBColor(0x19, 0x76, 0xD2)
LIGHT_BLUE = RGBColor(0xE3, 0xF2, 0xFD)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x21, 0x21, 0x21)
GRAY = RGBColor(0x61, 0x61, 0x61)
ORANGE = RGBColor(0xE6, 0x51, 0x00)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xC6, 0x28, 0x28)

def add_header_bar(slide, prs):
    """Add a dark blue header bar to the top of the slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, Inches(0.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

def add_footer(slide, prs, text="Gannamaneni Jaswanth — NITK Surathkal — M.Tech IT — 2026"):
    """Add footer text."""
    left = Inches(0.5)
    top = Inches(7.1)
    width = prs.slide_width - Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(8)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

def add_slide_title(slide, prs, title_text, subtitle_text=None):
    """Add title text on the header bar."""
    add_header_bar(slide, prs)
    txBox = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.08), prs.slide_width - Inches(1), Inches(0.55))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle_text:
        txBox2 = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.75), prs.slide_width - Inches(1), Inches(0.35))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle_text
        p2.font.size = Pt(12)
        p2.font.color.rgb = GRAY
        p2.font.italic = True

def add_body_text(slide, text_lines, left=0.6, top=1.2, width=11.5, font_size=14,
                  color=BLACK, bold_first=False, line_spacing=1.3):
    """Add body text with bullet points."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(4)
        if bold_first and i == 0:
            p.font.bold = True

def add_image_slide(slide, prs, title, image_path, subtitle=None,
                    img_left=0.3, img_top=1.0, img_width=12.5):
    """Add a slide with a title and centered image."""
    add_slide_title(slide, prs, title, subtitle)
    if os.path.exists(image_path):
        slide.shapes.add_picture(
            image_path, Inches(img_left), Inches(img_top),
            width=Inches(img_width))
    add_footer(slide, prs)

def add_two_col_text(slide, left_lines, right_lines, top=1.2):
    """Add two-column text layout."""
    # Left column
    txL = slide.shapes.add_textbox(
        Inches(0.5), Inches(top), Inches(5.5), Inches(5.5))
    tfL = txL.text_frame
    tfL.word_wrap = True
    for i, line in enumerate(left_lines):
        p = tfL.paragraphs[0] if i == 0 else tfL.add_paragraph()
        p.text = line
        p.font.size = Pt(13)
        p.font.color.rgb = BLACK
        p.space_after = Pt(3)

    # Right column
    txR = slide.shapes.add_textbox(
        Inches(6.5), Inches(top), Inches(5.5), Inches(5.5))
    tfR = txR.text_frame
    tfR.word_wrap = True
    for i, line in enumerate(right_lines):
        p = tfR.paragraphs[0] if i == 0 else tfR.add_paragraph()
        p.text = line
        p.font.size = Pt(13)
        p.font.color.rgb = BLACK
        p.space_after = Pt(3)


def build_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # blank layout

    # ══════════════════════════════════════════════════════════
    # SLIDE 1: TITLE
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    # Full blue background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()

    # Title
    tx = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.3), Inches(1.5))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Imitating Fund Manager Decisions"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "A Causally-Informed Knowledge Graph Approach\nfor Portfolio Construction"
    p2.font.size = Pt(22)
    p2.font.color.rgb = RGBColor(0xBB, 0xDE, 0xFB)
    p2.alignment = PP_ALIGN.CENTER

    # Author info
    tx2 = slide.shapes.add_textbox(Inches(2), Inches(4.2), Inches(9.3), Inches(2.5))
    tf2 = tx2.text_frame
    tf2.word_wrap = True
    lines = [
        ("Gannamaneni Jaswanth", 20, True),
        ("M.Tech — Information Technology", 16, False),
        ("", 10, False),
        ("Guide: Dr. Biju R. Mohan", 16, False),
        ("", 10, False),
        ("Department of Information Technology", 14, False),
        ("National Institute of Technology Karnataka, Surathkal", 14, False),
        ("2025–2026", 14, False),
    ]
    for i, (text, size, bold) in enumerate(lines):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # ══════════════════════════════════════════════════════════
    # SLIDE 2: OUTLINE
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_slide_title(slide, prs, "Presentation Outline")
    add_two_col_text(slide,
        ["1.  Problem Statement & Motivation",
         "2.  Literature Review & Research Gaps",
         "3.  Research Objectives",
         "4.  System Architecture",
         "5.  Phase 1: Data Pipeline",
         "6.  Phase 1: Causal Discovery",
         "7.  Phase 1: KG Construction",
         "8.  Phase 1: Novel Evaluation Metrics",
         "9.  Phase 1: Results"],
        ["10. Phase 2: CBR Inference Engine",
         "11. Phase 2: Causal Path Transformer",
         "12. Phase 2: HGT & CI-HGT (Novel)",
         "13. Phase 2: 7-Way Comparison",
         "14. Phase 2: Walk-Forward Backtest",
         "15. Phase 2: Ablation & XAI",
         "16. Novel Contributions",
         "17. Limitations & Future Work",
         "18. Conclusion & References"],
        top=1.2)
    add_footer(slide, prs)

    # ══════════════════════════════════════════════════════════
    # SLIDE 3: PROBLEM STATEMENT
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_slide_title(slide, prs, "Problem Statement")
    add_body_text(slide, [
        "■  Fund managers make complex portfolio decisions (BUY / HOLD / SELL)",
        "    based on years of experience, intuition, and multi-factor analysis.",
        "",
        "■  These decisions are rarely documented in structured, machine-readable form.",
        "    The knowledge resides as implicit expertise — difficult to transfer or replicate.",
        "",
        "■  When a fund manager leaves, the investment strategy is effectively lost.",
        "    New managers must rebuild decision frameworks from scratch.",
        "",
        "■  Problem: Can we capture, structure, and replicate the decision-making",
        "    patterns of experienced fund managers using observable data?",
        "",
        "■  Challenge: Standard ML (e.g., XGBoost on tabular features) treats decisions",
        "    as isolated predictions. It ignores causal relationships, temporal structure,",
        "    and the interconnected nature of financial markets.",
    ], font_size=13)
    add_footer(slide, prs)

    # ══════════════════════════════════════════════════════════
    # SLIDE 4: MOTIVATION
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_slide_title(slide, prs, "Motivation",
                    "Why Knowledge Graphs + Causal Inference?")
    add_two_col_text(slide,
        ["Why Knowledge Graphs?",
         "",
         "• Capture relational structure: funds,",
         "  stocks, sectors, time, regimes",
         "• Represent causal dependencies",
         "  as first-class graph edges",
         "• Enable graph-based reasoning",
         "  (CBR, path queries, GNNs)",
         "• Natural fit for heterogeneous",
         "  financial data",
         "",
         "(Hogan et al., 2021)"],
        ["Why Causal Inference?",
         "",
         "• Correlation ≠ Causation —",
         "  spurious features mislead models",
         "• Causal features are more stable",
         "  across market regimes",
         "• Three complementary methods:",
         "  Granger, ICP, DML — each captures",
         "  different aspects of causality",
         "• Enables interpretable explanations",
         "  (\"pct_nav causes SELL decisions\")",
         "",
         "(Peters et al., 2016; Chernozhukov et al., 2018)"],
        top=1.2)
    add_footer(slide, prs)

    # ══════════════════════════════════════════════════════════
    # SLIDE 5: LITERATURE REVIEW
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_slide_title(slide, prs, "Literature Review")
    add_body_text(slide, [
        "Knowledge Graphs in Finance",
        "  • KGs used for fraud detection, risk assessment, and entity linking",
        "    (Hogan et al., 2021). Limited application to portfolio decision mimicry.",
        "",
        "Causal Inference in Financial Markets",
        "  • Granger causality widely used for time-series analysis (Granger, 1969).",
        "  • ICP provides invariance-based causal discovery (Peters et al., 2016).",
        "  • DML enables debiased treatment effect estimation (Chernozhukov et al., 2018).",
        "  • No prior work combines all three methods into a unified KG.",
        "",
        "Graph Neural Networks for Finance",
        "  • GNNs applied to stock prediction (R-GCN: Schlichtkrull et al., 2018).",
        "  • HGT handles heterogeneous graphs (Hu et al., 2020).",
        "  • No prior work applies HGT to fund manager decision mimicry.",
        "",
        "Case-Based Reasoning",
        "  • CBR retrieves similar past cases for decision support (Aamodt & Plaza, 1994).",
        "  • Graph kernel similarity (WL kernel) enables structural matching",
        "    (Shervashidze et al., 2011).",
    ], font_size=12)
    add_footer(slide, prs)

    # ══════════════════════════════════════════════════════════
    # SLIDE 6: RESEARCH GAPS
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_slide_title(slide, prs, "Research Gaps & Contributions")
    add_two_col_text(slide,
        ["Identified Gaps:",
         "",
         "1. No system uses causally-informed",
         "   KGs to imitate fund managers",
         "",
         "2. No unified framework combines",
         "   Granger + ICP + DML for KG",
         "   enrichment in finance",
         "",
         "3. No standard metrics exist to",
         "   evaluate causally-informed KGs",
         "",
         "4. HGT has not been applied with",
         "   causal edge modulation (CI-HGT)"],
        ["Our Contributions:",
         "",
         "1. End-to-end 16-step pipeline:",
         "   data → KG → mimicry → backtest",
         "",
         "2. Tri-method causal discovery",
         "   with KG integration",
         "",
         "3. Three novel evaluation metrics:",
         "   CSCS, SCSI, DMF",
         "",
         "4. CI-HGT: CausalGate mechanism",
         "   that modulates GNN message",
         "   passing via causal edge strengths"],
        top=1.2)
    add_footer(slide, prs)

    # ══════════════════════════════════════════════════════════
    # SLIDE 7: RESEARCH OBJECTIVES
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_slide_title(slide, prs, "Research Objectives")
    add_body_text(slide, [
        "Primary Objective:",
        "  To design and develop a system that analyzes patterns in stock selection",
        "  and allocation from the experience of previous fund manager decisions and",
        "  uses this learning to imitate the decision-making process of a fund manager",
        "  for future investment scenarios.",
        "",
        "Specific Objectives:",
        "",
        "O1.  Construct a temporally-structured Knowledge Graph from Indian mutual fund",
        "     portfolio data (32 funds × 46 months × 1,057 stocks).",
        "",
        "O2.  Enrich the KG with multi-method causal discovery (Panel Granger Causality,",
        "     Invariant Causal Prediction, Double Machine Learning).",
        "",
        "O3.  Develop KG-native mimicry models (CBR, Path Transformer, HGT, CI-HGT)",
        "     that operate directly on the graph structure.",
        "",
        "O4.  Evaluate mimicry fidelity (accuracy, F1, Cohen's κ) and financial",
        "     performance (Sharpe, Sortino, Information Ratio) through walk-forward backtest.",
    ], font_size=12.5)
    add_footer(slide, prs)

    # ══════════════════════════════════════════════════════════
    # SLIDE 8: SYSTEM ARCHITECTURE (Image)
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_image_slide(slide, prs,
        "System Architecture — End-to-End Pipeline",
        os.path.join(DIAG_DIR, '01_methodology_flow.png'),
        img_left=0.2, img_top=0.85, img_width=12.9)

    # ══════════════════════════════════════════════════════════
    # SLIDE 9: DATA PIPELINE (Image)
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_image_slide(slide, prs,
        "Phase 1A: Data Acquisition & Feature Engineering",
        os.path.join(DIAG_DIR, '10_data_pipeline.png'),
        subtitle="Steps 00–08: Five data sources → unified panel dataset",
        img_left=0.2, img_top=1.1, img_width=12.9)

    # ══════════════════════════════════════════════════════════
    # SLIDE 10: DATA DETAILS
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_slide_title(slide, prs, "Data Sources & Feature Engineering",
                    "Steps 00–08")
    add_two_col_text(slide,
        ["Data Sources:",
         "  • AMFI India: Monthly portfolio disclosures",
         "    for 32 open-ended equity mutual funds",
         "  • Screener.in: Quarterly fundamentals",
         "    (EPS, P/E, revenue, debt ratios)",
         "  • Kite/NSE API: Daily OHLCV prices",
         "    aggregated to monthly bars",
         "  • FinBERT (Araci, 2019): NLP-based",
         "    news sentiment per stock per month",
         "  • RBI/FRED: Macro indicators (GDP,",
         "    VIX, crude oil, repo rate, US 10Y)"],
        ["Feature Engineering (Step 08):",
         "  • Lag features: 1–6 month lags",
         "    for all numeric variables",
         "  • Momentum: 3/6/12-month returns",
         "  • Rolling statistics: mean, std",
         "  • Cross-sectional ranks within month",
         "  • Pruning: remove correlated > 0.95",
         "  • Final: 83,643 observations × 58 features",
         "",
         "Output Datasets:",
         "  • CAUSAL_DISCOVERY_DATASET.csv (raw)",
         "  • LPCMCI_READY.csv (standardized)"],
        top=1.2)
    add_footer(slide, prs)

    # ══════════════════════════════════════════════════════════
    # SLIDE 11: CAUSAL DISCOVERY FRAMEWORK (Image)
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_image_slide(slide, prs,
        "Phase 1B: Causal Discovery Framework",
        os.path.join(DIAG_DIR, '03_causal_discovery_framework.png'),
        subtitle="Three complementary methods: Granger + ICP + DML",
        img_left=0.2, img_top=1.1, img_width=12.9)

    # ══════════════════════════════════════════════════════════
    # SLIDE 12: GRANGER CAUSALITY
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_slide_title(slide, prs, "Panel Granger Causality (Step 09)",
                    "Granger, 1969; adapted for panel data with entity fixed effects")
    add_body_text(slide, [
        "Method:",
        "  • Panel OLS regression with entity (fund×stock) fixed effects",
        "  • Tests: does variable X at lag t-k predict action_ordinal at t?",
        "  • Lags tested: 1–6 months",
        "  • FDR correction (Benjamini-Hochberg) at α = 0.05",
        "  • Stratified analysis by fund type",
        "",
        "Results:",
        "  • 106 statistically significant causal edges (100% FDR-significant)",
        "  • 10 unique causal variables identified",
        "  • Average lag: 3.29 months (mid-term effects dominate)",
        "  • Lag distribution: fairly uniform across 1–6 months",
        "  • Strongest predictor: pct_nav (portfolio allocation percentage)",
        "",
        "KG Integration:",
        "  • Each significant relationship → GRANGER_CAUSES edge in Neo4j",
        "  • Properties: beta, p_fdr, partial_r², lag, direction",
        "",
        "Reference: Granger, C. W. J. (1969). Investigating causal relations by econometric",
        "models and cross-spectral methods. Econometrica, 37(3), 424–438.",
    ], font_size=12)
    add_footer(slide, prs)

    # ══════════════════════════════════════════════════════════
    # SLIDE 13: ICP
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_slide_title(slide, prs, "Invariant Causal Prediction — ICP (Step 09a)",
                    "Peters, Bühlmann, & Meinshausen, 2016")
    add_body_text(slide, [
        "Method (v7 — Grow-Shrink Markov Blanket + Exhaustive ICP):",
        "  • Step 1: Grow-Shrink algorithm discovers Markov Blanket for each target",
        "    (action_ordinal, is_buy, is_sell) — ~18-28 variables per target",
        "  • Step 2: Rich environments = quarter × regime × temporal_position",
        "    (e.g., Q1_crisis_early, Q3_bull_mid) — up to 24 environments",
        "  • Step 3: Exhaustive subset ICP (Peters et al., 2016) on MB variables",
        "  • Step 4: Soft intersection (80% threshold) across strata",
        "",
        "Results:",
        "  • Markov Blanket: 46 unique variables (union of 3 targets)",
        "  • ICP parents found for is_sell target (confidence ≥ 0.30):",
        "    pct_nav_lag4/5/6, sentiment_mean, avg_neutral_prob (5 parents)",
        "  • ICP parents for action_ordinal/is_buy: ZERO (invariance too strict",
        "    with only 4 pooled environments reaching sample threshold)",
        "",
        "KG Integration:",
        "  • ICP parents → CAUSES edges with confidence scores",
        "  • Markov Blanket used for M1 feature selection (46 features → 61.9% accuracy)",
        "",
        "Reference: Peters, J., Bühlmann, P., & Meinshausen, N. (2016). Causal inference",
        "by using invariant prediction. JRSS-B, 78(5), 947–1012.",
    ], font_size=11.5)
    add_footer(slide, prs)

    # ══════════════════════════════════════════════════════════
    # SLIDE 14: DML
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_slide_title(slide, prs, "Double Machine Learning — DML (Step 09b)",
                    "Chernozhukov et al., 2018")
    add_body_text(slide, [
        "Method:",
        "  • Estimates causal treatment effects while controlling for confounders",
        "  • Uses Random Forest for residualization (removing confounding)",
        "  • 5-fold cross-fitting to avoid overfitting bias",
        "  • Produces debiased θ̂ (theta-hat) estimates with 95% CI bounds",
        "",
        "Key Advantage over Granger:",
        "  • Granger tests predictive causality (temporal precedence)",
        "  • DML tests interventional causality (\"what happens if I change X?\")",
        "  • Robust to high-dimensional confounders via ML-based nuisance estimation",
        "",
        "Results (v4 — v5 auto-discovery pending re-run):",
        "  • 120 total effects estimated, 98 statistically significant",
        "  • 35 unique significant treatment variables",
        "  • Top effect: real_interest_rate → is_sell (θ̂ = 3.65)",
        "  • india_vix_close → action_ordinal (θ̂ = 0.48)",
        "",
        "KG Integration:",
        "  • Each significant DML effect → CAUSAL_EFFECT edge in Neo4j",
        "  • Properties: theta_hat, ci_lower_95, ci_upper_95, t_statistic, p_value",
        "",
        "Reference: Chernozhukov, V., et al. (2018). Double/debiased machine learning.",
        "The Econometrics Journal, 21(1), C1–C68.",
    ], font_size=11.5)
    add_footer(slide, prs)

    # ══════════════════════════════════════════════════════════
    # SLIDE 15: MULTI-METHOD CONSENSUS (Image)
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_image_slide(slide, prs,
        "Multi-Method Causal Consensus",
        os.path.join(DIAG_DIR, '12_multi_method_consensus.png'),
        subtitle="Granger (106), ICP (32 parents), DML (117 effects) — multi-method causal validation",
        img_left=1.5, img_top=1.0, img_width=10.0)

    # ══════════════════════════════════════════════════════════
    # SLIDE 16: KG CONSTRUCTION
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_slide_title(slide, prs, "Knowledge Graph Construction (Steps 10–11b)",
                    "Neo4j graph database — temporal + causal layers")
    add_body_text(slide, [
        "Temporal Layer (Step 10):",
        "  • Fund → HOLDS → Stock edges with temporal properties (pct_nav, tenure,",
        "    position_action, monthly_return, rank, consensus)",
        "  • Fund → EXITED → Stock edges for position exits",
        "  • TimePeriod → NEXT → TimePeriod chain (temporal ordering)",
        "  • MarketRegime nodes: VIX-based regime classification",
        "  • FundSnapshot / StockSnapshot: monthly state captures",
        "",
        "Causal Layer (Steps 11, 11b):",
        "  • GRANGER_CAUSES edges: 106 Panel Granger results",
        "  • CAUSES edges: ICP invariance results",
        "  • CAUSAL_EFFECT edges: DML treatment effects",
        "  • All causal edges carry statistical properties (p-values, effect sizes, CIs)",
        "",
        "Scale:",
        "  • 27,842 nodes (10 types) — 143,814 edges (13 types)",
        "  • Density: 0.000186 (sparse, as expected for real-world financial graph)",
        "  • Connected ratio: 99.99% (only 1 isolated node)",
        "  • Average degree: 10.33",
    ], font_size=12)
    add_footer(slide, prs)

    # ══════════════════════════════════════════════════════════
    # SLIDE 17: KG SCHEMA (Image)
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_image_slide(slide, prs,
        "Knowledge Graph Schema",
        os.path.join(DIAG_DIR, '02_kg_schema.png'),
        subtitle="10 node types, 13 relationship types — heterogeneous property graph",
        img_left=0.5, img_top=1.0, img_width=12.3)

    # ══════════════════════════════════════════════════════════
    # SLIDE 18: NOVEL METRICS (Image)
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_image_slide(slide, prs,
        "Novel KG Evaluation Metrics (Step 12b)",
        os.path.join(DIAG_DIR, '06_novel_evaluation_metrics.png'),
        subtitle="Proposed metrics — no prior standard exists for causally-informed KG evaluation",
        img_left=0.2, img_top=1.1, img_width=12.9)

    # ══════════════════════════════════════════════════════════
    # SLIDE 19: PHASE 1 RESULTS (Image)
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_image_slide(slide, prs,
        "Phase 1 Results Summary",
        os.path.join(DIAG_DIR, '07_phase1_results_table.png'),
        subtitle="KG structural statistics + causal discovery + novel evaluation metrics",
        img_left=0.2, img_top=1.0, img_width=12.9)

    # ══════════════════════════════════════════════════════════
    # SLIDE 20: COMPETENCY QUESTIONS (Image)
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_image_slide(slide, prs,
        "Inferential Utility — Competency Questions",
        os.path.join(DIAG_DIR, '13_competency_questions.png'),
        subtitle="8/10 domain questions answered by the KG (80% inferential utility)",
        img_left=0.3, img_top=1.0, img_width=12.5)

    # ══════════════════════════════════════════════════════════
    # SLIDE 21: ACTION DISTRIBUTION (Image)
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_image_slide(slide, prs,
        "Decision Distribution Analysis",
        os.path.join(DIAG_DIR, '11_action_distribution.png'),
        subtitle="Class imbalance: HOLD dominates (48.7%) — requires weighted evaluation metrics",
        img_left=1.5, img_top=1.0, img_width=10.0)

    # ══════════════════════════════════════════════════════════
    # SLIDE 22: PHASE 2 OVERVIEW (Image)
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_image_slide(slide, prs,
        "Phase 2: KG-Based Fund Manager Mimicry",
        os.path.join(DIAG_DIR, '04_phase2_architecture.png'),
        subtitle="Three KG-native models → 7-way comparison → walk-forward backtest",
        img_left=0.2, img_top=1.0, img_width=12.9)

    # ══════════════════════════════════════════════════════════
    # SLIDE 23: CBR ENGINE
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_slide_title(slide, prs, "CBR Inference Engine (Step 13)",
                    "Case-Based Reasoning on KG Subgraphs — Aamodt & Plaza, 1994")
    add_body_text(slide, [
        "Approach:",
        "  • For each new decision (fund × stock × month), extract a local KG subgraph",
        "  • Subgraph includes: HOLDS properties, causal context, sector, regime",
        "  • Compute Weisfeiler-Lehman graph kernel similarity with historical cases",
        "    (Shervashidze et al., 2011)",
        "  • k-NN retrieval: find most similar past decisions and their outcomes",
        "  • Majority vote among k nearest neighbors → predicted action",
        "",
        "Walk-Forward Cross-Validation:",
        "  • Train on months 1..t, predict month t+1",
        "  • Embargo: 1 month gap to prevent data leakage",
        "  • Each fold produces per-decision predictions",
        "",
        "Output:",
        "  • cbr_decision_predictions.csv: (Fund, ISIN, month, predicted, actual, confidence)",
        "  • Feeds into step14b (comparison) and step16 (backtest)",
        "",
        "Key: This model operates DIRECTLY on the KG structure —",
        "it cannot function without the Knowledge Graph.",
    ], font_size=12)
    add_footer(slide, prs)

    # ══════════════════════════════════════════════════════════
    # SLIDE 24: PATH TRANSFORMER
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_slide_title(slide, prs, "Causal Path Transformer (Step 13a)",
                    "Transformer attention over KG causal paths — Vaswani et al., 2017")
    add_body_text(slide, [
        "Approach:",
        "  • Extract causal paths from KG: sequences of nodes connected by causal edges",
        "  • Encode paths using positional + type embeddings",
        "  • Apply Transformer self-attention to learn which paths matter for each decision",
        "  • d_model=128, 8 attention heads, 3 encoder layers",
        "",
        "Training:",
        "  • Warmup (10% of epochs) + cosine annealing learning rate schedule",
        "  • Walk-forward temporal split (no future data leakage)",
        "  • Fallback: sklearn GradientBoostingClassifier if PyTorch unavailable",
        "",
        "Novelty:",
        "  • Learns path-level importance: not just which features matter,",
        "    but which SEQUENCES of causal relationships predict fund manager decisions",
        "  • Example path: pct_nav → GRANGER_CAUSES → rsi → CAUSES → action_ordinal",
        "",
        "Output:",
        "  • path_decision_predictions.csv with per-decision probabilities",
        "  • path_embeddings.npy for downstream style clustering (step 16c)",
    ], font_size=12)
    add_footer(slide, prs)

    # ══════════════════════════════════════════════════════════
    # SLIDE 25: CI-HGT (Image)
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_image_slide(slide, prs,
        "CI-HGT: Causally-Informed Heterogeneous Graph Transformer (Novel)",
        os.path.join(DIAG_DIR, '05_cihgt_architecture.png'),
        subtitle="Novel contribution: CausalGate modulates message passing via causal edge strengths",
        img_left=0.2, img_top=1.0, img_width=12.9)

    # ══════════════════════════════════════════════════════════
    # SLIDE 26: HGT vs CI-HGT DETAILS
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_slide_title(slide, prs, "HGT vs CI-HGT — Technical Details (Step 13b)",
                    "Hu et al., 2020 + our CausalGate extension")
    add_two_col_text(slide,
        ["HGT (Baseline):",
         "  • Heterogeneous Graph Transformer",
         "    (Hu et al., 2020)",
         "  • Type-specific attention for each",
         "    node type and edge type",
         "  • Hidden dim = 128, 4 heads, 3 layers",
         "  • Operates on full KG structure",
         "  • Edge classifier: concat(h_fund,",
         "    h_stock, edge_features) → MLP",
         "  • Predicts: P(BUY|HOLD|SELL)",
         "    per fund × stock × month"],
        ["CI-HGT (Novel Extension):",
         "  • Adds CausalGate module after HGT",
         "  • gate = σ(W · causal_edge_attr)",
         "  • h_src_gated = h_src ⊙ gate",
         "  • Scatter-mean aggregation to targets",
         "  • 0.1-scaled residual update",
         "",
         "  Key Insight:",
         "  • Strong causal evidence (high ICP",
         "    confidence, low Granger p-value)",
         "    is amplified in message passing",
         "  • Weak evidence is attenuated",
         "  • The gate LEARNS the threshold"],
        top=1.2)
    add_footer(slide, prs)

    # ══════════════════════════════════════════════════════════
    # SLIDE 27: 7-WAY COMPARISON (Image)
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_image_slide(slide, prs,
        "7-Way Model Comparison (Step 14b)",
        os.path.join(DIAG_DIR, '08_phase2_comparison_table.png'),
        subtitle="Cohen's κ measures agreement with fund manager (Cohen, 1960) — Walk-forward evaluation",
        img_left=0.2, img_top=1.2, img_width=12.9)

    # ══════════════════════════════════════════════════════════
    # SLIDE 28: BACKTEST (Image)
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_image_slide(slide, prs,
        "Walk-Forward Backtest with Real Returns (Step 16)",
        os.path.join(DIAG_DIR, '09_backtest_results_table.png'),
        subtitle="Sharpe (1966), Sortino (1991), Information Ratio (Grinold & Kahn, 2000)",
        img_left=0.2, img_top=1.2, img_width=12.9)

    # ══════════════════════════════════════════════════════════
    # SLIDE 29: XAI
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_slide_title(slide, prs, "Explainable AI — KG-Grounded Explanations (Step 15)",
                    "Lundberg & Lee, 2017 — SHAP values grounded in KG structure")
    add_body_text(slide, [
        "Approach:",
        "  • SHAP (SHapley Additive exPlanations) feature importance",
        "    (Lundberg & Lee, 2017)",
        "  • Ground SHAP explanations in KG: map important features to",
        "    causal edges in the Knowledge Graph",
        "  • CI-based counterfactual confidence: use DML confidence intervals",
        "    to estimate robustness of each explanation",
        "",
        "Metrics:",
        "  • Faithfulness: do explanations match model behavior?",
        "  • avg_ci_faithfulness: faithfulness weighted by causal CI width",
        "  • KG grounding ratio: what fraction of top features have KG backing?",
        "",
        "Example Explanation:",
        "  \"The model predicted BUY because pct_nav = 2.3% (SHAP = +0.15),",
        "   which is causally linked to BUY decisions via GRANGER_CAUSES",
        "   (β = 0.045, p_fdr < 0.001, lag = 2 months).\"",
        "",
        "This bridges the gap between black-box ML and interpretable KG reasoning.",
    ], font_size=12)
    add_footer(slide, prs)

    # ══════════════════════════════════════════════════════════
    # SLIDE 30: ABLATION STUDY
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_slide_title(slide, prs, "Ablation Study Framework (Step 16b)",
                    "Essential for thesis defense — proves each component's contribution")
    add_body_text(slide, [
        "1. Feature Group Ablation:",
        "   • Remove each feature group (technical, fundamental, macro, sentiment,",
        "     position, causal) one at a time → measure accuracy/F1 drop",
        "   • Proves: each feature group contributes meaningfully",
        "",
        "2. Causal Method Ablation:",
        "   • Granger-only vs ICP-only vs DML-only vs intersection vs union",
        "   • Proves: three methods are complementary, not redundant",
        "",
        "3. Model Architecture Ablation:",
        "   • LogReg vs RF vs GBM vs XGBoost vs CBR vs PathTransformer vs HGT vs CI-HGT",
        "   • Proves: KG-native models add value over tabular ML",
        "",
        "4. ICP Confidence Threshold Ablation:",
        "   • Threshold variants: 0.10 / 0.25 / 0.50",
        "   • Shows sensitivity of causal feature selection to confidence cutoff",
        "",
        "Statistical Rigor:",
        "   • Mean ± std across walk-forward folds",
        "   • Paired t-test between full model and each ablation variant",
        "   • Cohen's d effect size for practical significance",
    ], font_size=12)
    add_footer(slide, prs)

    # ══════════════════════════════════════════════════════════
    # SLIDE 31: EVALUATION METRICS SUMMARY
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_slide_title(slide, prs, "Evaluation Metrics — Complete Framework")
    add_two_col_text(slide,
        ["Mimicry Metrics (Phase 2):",
         "  • Accuracy: overall correct predictions",
         "  • Weighted F1: handles class imbalance",
         "    (van Rijsbergen, 1979)",
         "  • Cohen's κ: inter-rater agreement",
         "    beyond chance (Cohen, 1960)",
         "  • Decision Agreement: % matching",
         "    fund manager's actual decisions",
         "  • BUY/SELL Recall: per-class sensitivity",
         "",
         "KG Quality Metrics (Phase 1):",
         "  • CSCS: Causal-Semantic Coherence",
         "  • SCSI: Stratified Causal Stability",
         "  • DMF: Decision-Mimicry Faithfulness",
         "  • Inferential Utility (competency Qs)"],
        ["Financial Metrics (Backtest):",
         "  • Sharpe Ratio: risk-adjusted return",
         "    (Sharpe, 1966)",
         "  • Sortino Ratio: downside risk only",
         "    (Sortino & van der Meer, 1991)",
         "  • Information Ratio: active return per",
         "    unit of tracking error",
         "    (Grinold & Kahn, 2000)",
         "  • Max Drawdown: peak-to-trough loss",
         "  • Calmar Ratio: return / max drawdown",
         "  • Hit Rate: % of positive months",
         "  • Turnover: portfolio churn rate",
         "",
         "Statistical Tests:",
         "  • Paired t-test for ablation",
         "  • Cohen's d effect size"],
        top=1.2)
    add_footer(slide, prs)

    # ══════════════════════════════════════════════════════════
    # SLIDE 32: NOVEL CONTRIBUTIONS
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_slide_title(slide, prs, "Novel Contributions — Summary")
    add_body_text(slide, [
        "1. Causally-Informed Knowledge Graph for Fund Manager Mimicry",
        "   • First system to combine Granger + ICP + DML into a unified KG",
        "     for portfolio decision imitation (no prior work)",
        "",
        "2. CI-HGT: CausalGate Mechanism",
        "   • Novel extension of HGT (Hu et al., 2020) where causal edge strengths",
        "     modulate GNN message passing via a learned gating function",
        "   • Amplifies strong causal evidence, attenuates weak evidence",
        "",
        "3. Three Novel KG Evaluation Metrics",
        "   • CSCS: measures theory alignment of causal edges",
        "   • SCSI: measures cross-segment stability of causal structure",
        "   • DMF: measures KG-classifier feature alignment",
        "   • No prior standard exists for evaluating causally-informed KGs",
        "",
        "4. End-to-End Reproducible Pipeline",
        "   • 16-step pipeline from raw data to backtested returns",
        "   • Every step is a standalone Python script with clear I/O",
        "   • Walk-forward temporal evaluation prevents future data leakage",
    ], font_size=12.5)
    add_footer(slide, prs)

    # ══════════════════════════════════════════════════════════
    # SLIDE 33: LIMITATIONS
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_slide_title(slide, prs, "Limitations")
    add_body_text(slide, [
        "Data Limitations:",
        "  • Monthly portfolio disclosures only (no daily/intraday granularity)",
        "  • 32 funds × 46 months — larger universe would improve generalizability",
        "  • Indian mutual funds only — results may not transfer to other markets",
        "",
        "Methodological Limitations:",
        "  • Granger causality assumes linear relationships",
        "  • ICP finds parents only for is_sell (not action_ordinal/is_buy) —",
        "    multi-class targets create invariance challenges across environments",
        "  • DML effect sizes are heterogeneous (real_interest_rate θ̂=3.65 vs",
        "    most features θ̂ < 0.10) — some may not be economically significant",
        "",
        "Model Limitations:",
        "  • KG models (HGT/CI-HGT) achieve lower accuracy (48.6%) but better",
        "    risk-adjusted returns (Calmar=50) — accuracy ≠ financial performance",
        "  • CBR and PathTransformer have negative Sharpe ratios (-0.53, -0.22)",
        "  • CI-HGT CausalGate uses simple sigmoid gating — more complex",
        "    attention mechanisms could be explored",
        "  • M1 (Markov Blanket) needs re-run with corrected file paths",
    ], font_size=12)
    add_footer(slide, prs)

    # ══════════════════════════════════════════════════════════
    # SLIDE 34: FUTURE WORK
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_slide_title(slide, prs, "Future Work")
    add_body_text(slide, [
        "Short-term Extensions:",
        "  • Expand to 100+ funds and longer time horizon (5+ years)",
        "  • Add international fund data (US, EU) for cross-market validation",
        "  • Incorporate daily/weekly granularity for more temporal resolution",
        "",
        "Methodological Extensions:",
        "  • Non-linear causal discovery (PCMCI+, NOTEARS, DAG-GNN)",
        "  • Dynamic KG that updates causally as new data arrives (incremental)",
        "  • Reinforcement learning for portfolio optimization using KG as state space",
        "",
        "Architecture Extensions:",
        "  • Multi-task learning: predict action + allocation size simultaneously",
        "  • Graph Attention Networks v2 (GATv2) as alternative to HGT",
        "  • Temporal graph networks (TGN) for continuous-time modeling",
        "  • Large Language Model integration for natural language explanations",
        "",
        "Deployment:",
        "  • Real-time decision support system for fund managers",
        "  • API-based KG query interface for portfolio analytics",
    ], font_size=12)
    add_footer(slide, prs)

    # ══════════════════════════════════════════════════════════
    # SLIDE 35: CONCLUSION
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_slide_title(slide, prs, "Conclusion")
    add_body_text(slide, [
        "We presented an end-to-end system for imitating fund manager decisions",
        "using a causally-informed Knowledge Graph approach.",
        "",
        "Key Achievements:",
        "",
        "  ✓  Constructed a heterogeneous KG with 27,842 nodes and 143,814 edges",
        "     from 32 Indian mutual fund portfolios over 46 months",
        "",
        "  ✓  Enriched the KG with 282 causal edges from three complementary methods",
        "     (Granger: 106, ICP: 32, DML: 117) — multi-method validated",
        "",
        "  ✓  Developed KG-native mimicry models: HGT (F1=0.503, Sharpe=1.52),",
        "     CI-HGT (F1=0.500, Sharpe=1.50), KG Ensemble (Sharpe=2.45, 70.9% return)",
        "",
        "  ✓  Proposed three novel evaluation metrics (CSCS = 0.549, SCSI = 0.500,",
        "     DMF = 0.600) for causally-informed KG assessment",
        "",
        "  ✓  KG models achieve exceptional risk-adjusted returns: HGT Calmar=50.16,",
        "     CI-HGT Calmar=49.18 (minimal drawdown despite lower accuracy)",
        "",
        "  ✓  Walk-forward backtest validates: KG Ensemble (Sharpe=2.45) competitive",
        "     with M0 all-features (Sharpe=2.38) using only graph structure",
    ], font_size=12.5)
    add_footer(slide, prs)

    # ══════════════════════════════════════════════════════════
    # SLIDE 36-37: REFERENCES (APA)
    # ══════════════════════════════════════════════════════════
    refs_page1 = [
        "Aamodt, A., & Plaza, E. (1994). Case-based reasoning: Foundational issues, methodological variations, and system approaches. AI Communications, 7(1), 39–59.",
        "",
        "Araci, D. (2019). FinBERT: Financial sentiment analysis with pre-trained language models. arXiv preprint arXiv:1908.10063.",
        "",
        "Chen, T., & Guestrin, C. (2016). XGBoost: A scalable tree boosting system. Proceedings of the 22nd ACM SIGKDD, 785–794.",
        "",
        "Chernozhukov, V., Chetverikov, D., Demirer, M., Duflo, E., Hansen, C., Newey, W., & Robins, J. (2018). Double/debiased machine learning for treatment and structural parameters. The Econometrics Journal, 21(1), C1–C68.",
        "",
        "Cohen, J. (1960). A coefficient of agreement for nominal scales. Educational and Psychological Measurement, 20(1), 37–46.",
        "",
        "Fey, M., & Lenssen, J. E. (2019). Fast graph representation learning with PyTorch Geometric. ICLR Workshop on Representation Learning on Graphs and Manifolds.",
        "",
        "Granger, C. W. J. (1969). Investigating causal relations by econometric models and cross-spectral methods. Econometrica, 37(3), 424–438.",
        "",
        "Grinold, R. C., & Kahn, R. N. (2000). Active portfolio management (2nd ed.). McGraw-Hill.",
    ]

    refs_page2 = [
        "Hogan, A., Blomqvist, E., Cochez, M., d'Amato, C., Melo, G. D., Gutierrez, C., ... & Zimmermann, A. (2021). Knowledge graphs. ACM Computing Surveys, 54(4), 1–37.",
        "",
        "Hu, Z., Dong, Y., Wang, K., & Sun, Y. (2020). Heterogeneous graph transformer. Proceedings of The Web Conference 2020, 2704–2710.",
        "",
        "Lundberg, S. M., & Lee, S. I. (2017). A unified approach to interpreting model predictions. Advances in Neural Information Processing Systems, 30.",
        "",
        "Peters, J., Bühlmann, P., & Meinshausen, N. (2016). Causal inference by using invariant prediction: Identification and confidence intervals. Journal of the Royal Statistical Society: Series B, 78(5), 947–1012.",
        "",
        "Schlichtkrull, M., Kipf, T. N., Bloem, P., Van Den Berg, R., Titov, I., & Welling, M. (2018). Modeling relational data with graph convolutional networks. European Semantic Web Conference, 593–607.",
        "",
        "Sharpe, W. F. (1966). Mutual fund performance. The Journal of Business, 39(1), 119–138.",
        "",
        "Shervashidze, N., Schweitzer, P., Van Leeuwen, E. J., Mehlhorn, K., & Borgwardt, K. M. (2011). Weisfeiler-Lehman graph kernels. Journal of Machine Learning Research, 12, 2539–2561.",
        "",
        "Sortino, F. A., & van der Meer, R. (1991). Downside risk. The Journal of Portfolio Management, 17(4), 27–31.",
        "",
        "Vaswani, A., Shazeer, N., Parmar, N., Uszkoreit, J., Jones, L., Gomez, A. N., ... & Polosukhin, I. (2017). Attention is all you need. Advances in Neural Information Processing Systems, 30.",
    ]

    for refs, page in [(refs_page1, 1), (refs_page2, 2)]:
        slide = prs.slides.add_slide(blank_layout)
        add_slide_title(slide, prs, f"References ({page}/2)")
        add_body_text(slide, refs, font_size=10.5, top=1.0, left=0.6, width=12.0)
        add_footer(slide, prs)

    # ══════════════════════════════════════════════════════════
    # SLIDE 38: THANK YOU
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(9.3), Inches(3))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "\nQuestions & Discussion"
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(0xBB, 0xDE, 0xFB)
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = "\n\nGannamaneni Jaswanth\ngannamaneni.jaswanth@nitk.edu.in"
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(0x90, 0xCA, 0xF9)
    p3.alignment = PP_ALIGN.CENTER

    # Save
    prs.save(OUT_PPT)
    print(f"\n  PPT saved: {OUT_PPT}")
    print(f"  Total slides: {len(prs.slides)}")


if __name__ == '__main__':
    print("=" * 60)
    print("Generating Thesis Defense PPT")
    print("=" * 60)
    build_ppt()
    print("  Done!")
